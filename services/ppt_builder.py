from pptx import Presentation
from pathlib import Path
import datetime

from services.bible import parse_range, scroll_bible

# 템플릿 기준 master/layout 인덱스
M0 = 0
M1 = 1

# master0 layouts
L0_GREETINGS     = 0
L0_PREVIEW       = 1   # (ph 10,11,12)
L0_CONTACTUS     = 2
L0_WHITEPRAYER   = 5
L0_APOSTLES      = 6
L0_APOSTLES_1    = 7
L0_APOSTLES_2    = 8
L0_OFFERINGPRAY  = 9
L0_WHITEANNOUNCE = 10
L0_SERMON_TITLE  = 11  # (ph 10,11,12)
L0_BIBLE_PHRASE  = 12  # (ph 10,11)

# master1 layouts
L1_WORSHIP_BG    = 0
L1_WORSHIP_INTRO = 1
L1_WORSHIP_TITLE = 2   # (ph 10,11)
L1_WORSHIP_LYR   = 3   # (ph 10)
L1_BLACKPRAYER   = 4
L1_OFFERING      = 5
L1_LordsPrayer   = 6
L1_LordsPrayer_1 = 7
L1_LordsPrayer_2 = 8
L1_PRAYER        = 9   # (ph 10)


def clear_all_slides(p: Presentation):
    sldIdLst = p.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.rId
        p.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def add_slide(p: Presentation, master_idx: int, layout_idx: int):
    layout = p.slide_masters[master_idx].slide_layouts[layout_idx]
    return p.slides.add_slide(layout)


def set_ph(slide, idx: int, text: str):
    slide.placeholders[idx].text = text


def clean_song_title(title: str) -> str:
    return title.split("_")[0] if "_" in title else title


def add_song_section(p: Presentation, label_en: str, song: dict):
    title = clean_song_title(song.get("title", ""))
    lines = (song.get("lyrics", "") or "").replace("\r\n", "\n").split("\n")

    # 제목 슬라이드
    s_title = add_slide(p, M1, L1_WORSHIP_TITLE)
    set_ph(s_title, 10, title)
    set_ph(s_title, 11, label_en)

    # ---- 핵심: 빈 줄 처리 규칙 반영 ----
    buf = []  # 현재 슬라이드에 넣을 가사 줄(최대 2줄)

    def flush_buf():
        """buf에 쌓인 1~2줄을 슬라이드로 내보냄"""
        nonlocal buf
        if not buf:
            return
        s_lyr = add_slide(p, M1, L1_WORSHIP_LYR)
        set_ph(s_lyr, 10, "\n".join(buf))
        buf = []

    prev_flushed_linecount = 0  # 직전에 만든 가사 슬라이드가 몇 줄이었는지(1 or 2)

    for raw in lines:
        line = raw.strip()

        # 1) 빈 줄이면
        if line == "":
            # (A) 직전에 "1줄짜리 슬라이드"를 막 만든 상황이면 → 빈 슬라이드 만들지 말고 그냥 무시
            if not buf and prev_flushed_linecount == 1:
                continue

            # (B) 현재 buf에 1줄이 쌓인 상태에서 빈 줄이 나오면:
            #     -> 그 1줄은 그대로 슬라이드로 만들고(1줄짜리),
            #     -> 그리고 "빈 슬라이드"는 만들지 않음(너 요구사항: 1줄+빈줄이면 넘어감)
            if len(buf) == 1:
                flush_buf()
                prev_flushed_linecount = 1
                continue

            # (C) buf가 비어있고(=이미 2줄 단위로 flush 되었거나) 그냥 빈 줄만 나온 경우:
            #     -> 여기서만 "빈 슬라이드" 생성
            if not buf:
                add_slide(p, M1, L1_WORSHIP_LYR)  # 빈 화면용
                prev_flushed_linecount = 0
            continue

        # 2) 일반 가사 줄이면 buf에 추가
        buf.append(line)

        # 2줄이 채워지면 슬라이드 생성
        if len(buf) == 2:
            flush_buf()
            prev_flushed_linecount = 2

    # 마지막 남은 1줄 처리
    if buf:
        flush_buf()
        prev_flushed_linecount = 1 if len(buf) == 1 else 2


def build_ppt(plan: dict, template_path: Path, bible_dir: Path, out_dir: Path) -> Path:
    p = Presentation(str(template_path))
    clear_all_slides(p)

    date = datetime.datetime.now().date()
    sermon_title = plan.get("sermon_title", "")
    sermon_phrases = plan.get("sermon_phrases", [])

    # [1] 프리뷰 + 사도신경
    add_slide(p, M0, L0_GREETINGS)

    s_prev = add_slide(p, M0, L0_PREVIEW)
    set_ph(s_prev, 10, date.strftime("%Y.%m.%d 주일 예배"))
    set_ph(s_prev, 11, sermon_title)
    set_ph(s_prev, 12, sermon_phrases[0] if sermon_phrases else "")

    add_slide(p, M0, L0_CONTACTUS)
    add_slide(p, M0, L0_WHITEPRAYER)
    add_slide(p, M0, L0_APOSTLES)
    add_slide(p, M0, L0_APOSTLES_1)
    add_slide(p, M0, L0_APOSTLES_2)

    # [2] 찬양
    add_slide(p, M1, L1_WORSHIP_BG)
    add_slide(p, M1, L1_WORSHIP_INTRO)

    for s in plan["songs"]["praise"]:
        add_song_section(p, "Praise", s)

    add_slide(p, M1, L1_BLACKPRAYER)

    s_pray = add_slide(p, M1, L1_PRAYER)
    set_ph(s_pray, 10, plan.get("prayer", "기도 | "))

    # [3] 헌금
    add_slide(p, M1, L1_OFFERING)
    if plan["songs"].get("offering"):
        add_song_section(p, "Offering", plan["songs"]["offering"])

    add_slide(p, M0, L0_OFFERINGPRAY)
    add_slide(p, M0, L0_WHITEANNOUNCE)

    # [4] 설교 + 본문
    if sermon_phrases:
        s_sermon = add_slide(p, M0, L0_SERMON_TITLE)
        set_ph(s_sermon, 10, "말씀 | 박성준 전도사님")
        set_ph(s_sermon, 11, sermon_title)
        set_ph(s_sermon, 12, sermon_phrases[0])

        for phrase in sermon_phrases:
            book, chapter, start_v, end_v = parse_range(phrase)
            if not book:
                continue
            verses = scroll_bible(Path(bible_dir), book, chapter, start_v, end_v)

            for i, vno in enumerate(range(start_v, end_v + 1)):
                s_b = add_slide(p, M0, L0_BIBLE_PHRASE)
                set_ph(s_b, 10, verses[i])
                set_ph(s_b, 11, f"{book} {chapter}장 {vno}절")
            s_sermon = add_slide(p, M0, L0_SERMON_TITLE)
            set_ph(s_sermon, 10, "말씀 | 박성준 전도사님")
            set_ph(s_sermon, 11, sermon_title)
            set_ph(s_sermon, 12, sermon_phrases[0])


    # [5] 설후찬
    if plan["songs"].get("closing"):
        add_song_section(p, "Closing", plan["songs"]["closing"])

    # 주기도문
    add_slide(p, M1, L1_LordsPrayer)
    add_slide(p, M1, L1_LordsPrayer_1)
    add_slide(p, M1, L1_LordsPrayer_2)

    # 저장 (out/ 아래에 생성)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    out_fp = out_dir / f"{date.strftime('%Y%m%d')}_thepureum_out.pptx"
    p.save(str(out_fp))
    return out_fp
